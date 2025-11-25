import json
import os
from typing import ClassVar, Dict, Any, List
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from crewai.tools import BaseTool


class CreatePPTFromJsonTool(BaseTool):
    name: str = "create_ppt_from_json_tool"
    description: str = (
        "Reads a JSON file with slides using flexible casing for keys "
        "(title/content/bullets/notes/style/subslides) "
        "and generates a PowerPoint presentation."
    )

    DEFAULT_THEME: ClassVar[Dict[str, Dict[str, Any]]] = {
        "title":    {"font_size": 46, "bold": True, "color": "#0B3D91", "alignment": "center"},
        "section":  {"font_size": 40, "bold": True, "color": "#00796B", "alignment": "center"},
        "topic":    {"font_size": 32, "bold": True, "color": "#C2185B", "alignment": "left"},
        "subtopic": {"font_size": 28, "bold": False, "color": "#222222", "alignment": "left"},
        "body":     {"font_size": 20, "bold": False, "color": "#222222", "alignment": "left"},
    }

    def _get(self, node: dict, key: str, default=None):
        # Case-insensitive key access
        return next((v for k, v in node.items() if k.lower() == key.lower()), default)

    def _rgb(self, hex_string: str):
        try:
            return RGBColor.from_string(hex_string.replace("#", ""))
        except:
            return RGBColor(0, 0, 0)

    def _apply_style(self, paragraph, style_key: str):
        style = self.DEFAULT_THEME.get(style_key.lower(), self.DEFAULT_THEME["body"])
        paragraph.font.size = Pt(style["font_size"])
        paragraph.font.bold = style["bold"]
        paragraph.font.color.rgb = self._rgb(style["color"])
        paragraph.alignment = {
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }.get(style.get("alignment", "left").lower(), PP_ALIGN.LEFT)

    def _create_slide(self, prs, node):
        style = self._get(node, "style", "body").lower()
        layout = 0 if style == "title" else 1
        slide = prs.slides.add_slide(prs.slide_layouts[layout])

        # Title
        title = self._get(node, "title")
        if title:
            slide.shapes.title.text = title
            self._apply_style(slide.shapes.title.text_frame.paragraphs[0], style)

        # Body Placeholder
        body = slide.placeholders[1].text_frame
        body.clear()

        # Content
        content = self._get(node, "content")
        if content:
            p = body.paragraphs[0]
            p.text = content
            self._apply_style(p, "body")

        # Bullets
        bullets = self._get(node, "bullets", [])
        if isinstance(bullets, list):
            for bullet in bullets:
                bp = body.add_paragraph()
                bp.text = bullet
                bp.level = 1
                self._apply_style(bp, "body")

        # Notes
        notes = self._get(node, "notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _walk(self, prs, slides: List[Dict[str, Any]]):
        for slide in slides:
            self._create_slide(prs, slide)
            subslides = self._get(slide, "subslides", [])
            if subslides:
                self._walk(prs, subslides)

    def _output_path(self, json_path):
        directory, name = os.path.split(json_path)
        no_ext = os.path.splitext(name)[0].replace(" ", "_")
        return os.path.join(directory, f"{no_ext}.pptx")

    def _run(self, json_path: str) -> Dict[str, Any]:
        try:
            if not os.path.exists(json_path):
                return {"success": False, "error": "JSON file not found"}

            with open(json_path, "r") as f:
                data = json.load(f)

            slides = self._get(data, "slides")
            if not slides:
                return {"success": False, "error": "Missing `slides` key in JSON"}

            prs = Presentation()
            self._walk(prs, slides)

            output_ppt = self._output_path(json_path)
            prs.save(output_ppt)

            return {
                "success": True,
                "message": "PPT generated successfully",
                "ppt_path": output_ppt,
            }

        except Exception as e:
            return {"success": False, "error": str(e)}


if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser()
    parser.add_argument("--json_path", required=True)
    args = parser.parse_args()
    tool = CreatePPTFromJsonTool()
    print(tool.run(json_path=args.json_path))
